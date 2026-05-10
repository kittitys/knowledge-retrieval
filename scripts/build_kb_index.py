#!/usr/bin/env python3
"""
Build BM25 index for knowledge-base projects (v2.5).

Scans actual source directory, extracts/caches text content, builds BM25 index.
No tag generation. The index is used solely for Phase A full-text search.

Usage:
    python .agents/skills/knowledge-retrieval/scripts/build_kb_index.py --project <folder_name>

Output:
    knowledge-base/<folder>/.bm25_index/index/  -- BM25 index files
    knowledge-base/<folder>/.bm25_index/file_mapping.json  -- file list for search
    knowledge-base/<folder>/.bm25_index/file_metadata.json  -- snapshot for incremental check
"""

import os
import sys
import json
import re
import argparse
from pathlib import Path
from collections import defaultdict

# ── Config ──────────────────────────────────────────────────────────
# Workspace root: prefer cwd (has AGENTS.md), fall back to file-path heuristic
if (Path.cwd() / 'AGENTS.md').exists():
    WORKSPACE = Path.cwd()
else:
    WORKSPACE = Path(__file__).resolve().parent.parent.parent.parent  # ../../.. to reach workspace from skill path
KB_ROOT = WORKSPACE / "knowledge-base"

TEXT_EXTENSIONS = {".md", ".txt", ".rst", ".org", ".csv"}
NON_TEXT_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}

# PDF extraction (cached for speed; fallback to empty if fails)
def _extract_pdf_text(pdf_path: Path) -> str:
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(str(pdf_path), maxpages=500)
        if text and len(text.strip()) > 100:
            return text
        # Try PyMuPDF as fallback for PDFs pdfminer can't handle
        try:
            import fitz
            doc = fitz.open(str(pdf_path))
            text = "".join(page.get_text() for page in doc)
            doc.close()
            return text
        except Exception:
            return ""
    except Exception:
        # Last-resort fallback: PyMuPDF
        try:
            import fitz
            doc = fitz.open(str(pdf_path))
            text = "".join(page.get_text() for page in doc)
            doc.close()
            return text
        except Exception:
            return ""


def _extract_pptx_text(pptx_path: Path) -> str:
    """Extract text from PPTX, including GroupShape children and notes."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        texts = []
        prs = Presentation(str(pptx_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                # Recursively handle GroupShape
                def _collect_text(s):
                    if hasattr(s, "text_frame"):
                        for para in s.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                    if hasattr(s, "notes_slide"):
                        try:
                            notes = s.notes_slide.notes_text_frame.text.strip()
                            if notes:
                                texts.append("[NOTES] " + notes)
                        except Exception:
                            pass
                    if hasattr(s, "shapes"):  # GroupShape
                        for child in s.shapes:
                            _collect_text(child)
                _collect_text(shape)
        prs = None
        return "\n".join(texts)
    except Exception:
        return ""

# ── Helpers ─────────────────────────────────────────────────────────

def echo(msg):
    try:
        print(msg)
    except UnicodeEncodeError:
        safe = msg.encode("ascii", errors="replace").decode("ascii")
        print(safe)


def parse_data_structure(md_path: Path) -> dict | None:
    """Parse data_structure.md, return actual_location."""
    if not md_path.exists():
        return None
    text = md_path.read_text(encoding="utf-8", errors="replace")
    loc_match = re.search(r'^## 实际位置\s*>\s*(.+?)$', text, re.MULTILINE)
    if not loc_match:
        loc_match = re.search(r'^## Actual Location\s*>\s*(.+?)$', text, re.MULTILINE)
    if loc_match:
        return {"actual_location": loc_match.group(1).strip()}
    return {"actual_location": None}


def scan_project_files(project_root: Path, cache_dir: Path) -> list[dict]:
    """Scan directory for indexable files, extracting content for BM25 indexing."""
    files = []
    for entry in sorted(project_root.rglob("*")):
        if not entry.is_file():
            continue
        rel = entry.relative_to(project_root)
        parts = rel.parts
        if any(p.startswith(".") for p in parts):
            continue
        if entry.suffix.lower() == ".lnk":
            continue
        name = entry.name
        # Skip known processing artifacts (no project-specific overrides)
        if any(name.startswith(p) for p in ['_extracted_', '_doubao_extract_']):
            continue
        if name.startswith('temp_'):
            continue
        if '.bm25_index' in str(rel):
            continue

        info = {
            "path": str(entry.resolve()),
            "rel_path": str(rel),
            "name": entry.name,
            "suffix": entry.suffix.lower(),
            "size": entry.stat().st_size,
            "mtime": entry.stat().st_mtime,
            "content": "",
        }

        suf = info["suffix"]
        if suf in TEXT_EXTENSIONS:
            try:
                info["content"] = entry.read_text(encoding="utf-8", errors="replace")
            except Exception:
                pass
        elif suf in NON_TEXT_EXTENSIONS:
            cache_file = cache_dir / f"{entry.stem}.txt"
            if cache_file.exists():
                try:
                    info["content"] = cache_file.read_text(encoding="utf-8", errors="replace")
                except Exception:
                    pass
            else:
                # Extract directly if no cache exists yet
                extracted = ""
                if suf == ".pdf":
                    extracted = _extract_pdf_text(entry)
                elif suf == ".pptx":
                    extracted = _extract_pptx_text(entry)
                elif suf == ".docx":
                    try:
                        import docx
                        doc = docx.Document(str(entry))
                        extracted = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                    except Exception:
                        pass
                if extracted:
                    info["content"] = extracted
                    # Write to cache for Phase B reuse
                    try:
                        cache_dir.mkdir(parents=True, exist_ok=True)
                        cache_file.write_text(extracted, encoding="utf-8")
                    except Exception:
                        pass

        files.append(info)
    return files


def build_index_for_project(project_name: str, kb_dir: Path) -> dict:
    """Build BM25 index for one project. Returns summary dict."""
    import bm25s

    md_path = kb_dir / "data_structure.md"
    info = parse_data_structure(md_path)
    if info is None:
        return {"project": project_name, "status": "skipped", "reason": "no data_structure.md"}

    index_dir = kb_dir / ".bm25_index"
    index_dir.mkdir(parents=True, exist_ok=True)
    meta_path = index_dir / "file_metadata.json"

    # Check incremental — compare file snapshots
    existing = {}
    if meta_path.exists():
        try:
            existing = json.loads(meta_path.read_text("utf-8")).get("file_snapshot", {})
        except Exception:
            existing = {}

    # Determine source directory
    loc = info.get("actual_location")
    project_root = Path(loc) if loc and os.path.exists(loc) else None

    # Scan files
    all_files = []
    if project_root and project_root.exists():
        cache_dir = kb_dir / "cache"
        all_files = scan_project_files(project_root, cache_dir)

    if not all_files:
        return {"project": project_name, "status": "no_files", "reason": "source directory not accessible"}

    # Snapshot comparison for incremental skip
    snapshot = {}
    for f in all_files:
        snapshot[f["rel_path"]] = {"mtime": f["mtime"], "size": f["size"]}
    if snapshot == existing:
        echo(f"  [=] Index up-to-date ({len(all_files)} files)")
        return {"project": project_name, "status": "uptodate", "count": len(all_files)}

    # Build BM25 index
    echo(f"  [.] Indexing {len(all_files)} files...")
    corpora, mapping = [], []
    for f in all_files:
        content = f.get("content", "").strip()
        if not content:
            content = f["name"]
        corpora.append(content)
        mapping.append({"path": f["path"], "rel_path": f["rel_path"], "name": f["name"]})

    tokens = bm25s.tokenize(corpora, stopwords=None)
    retriever = bm25s.BM25()
    retriever.index(tokens)
    retriever.save(str(index_dir), corpus=corpora)

    # Save file mapping (needed by search_kb.py)
    (index_dir / "file_mapping.json").write_text(
        json.dumps(mapping, ensure_ascii=False), encoding="utf-8"
    )

    # Save minimal metadata (file snapshot only, no tags)
    metadata = {
        "built": True,
        "file_count": len(corpora),
        "file_snapshot": snapshot,
    }
    meta_path.write_text(json.dumps(metadata, ensure_ascii=False), encoding="utf-8")

    echo(f"  [OK] Indexed {len(corpora)}/{len(all_files)} files")
    return {"project": project_name, "status": "indexed", "indexed": len(corpora), "total": len(all_files)}


def main():
    parser = argparse.ArgumentParser(description="Build BM25 index for knowledge-base")
    parser.add_argument("--project", help="Specific project only")
    args = parser.parse_args()

    if not KB_ROOT.exists():
        echo(f"[!!] KB root not found: {KB_ROOT}")
        sys.exit(1)

    targets = [KB_ROOT / args.project] if args.project else sorted(KB_ROOT.iterdir())
    results = []

    for d in targets:
        if not d.is_dir() or d.name == "assets":
            continue
        echo(f"\n--- {d.name} ---")
        results.append(build_index_for_project(d.name, d))

    echo("\n== Summary ==")
    for r in results:
        st = r.get("status", "?")
        echo(f"  {r['project']}: {st}")


if __name__ == "__main__":
    main()

